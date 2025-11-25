"""Text extraction utilities for UFDR evidence files (F3.1)."""

import csv
import email
import json
import logging
import shutil
from contextlib import contextmanager
from dataclasses import dataclass
from pathlib import Path
from tempfile import NamedTemporaryFile
from typing import BinaryIO, Iterator

logger = logging.getLogger(__name__)


class MissingDependencyError(RuntimeError):
    """Raised when an optional dependency required for extraction is missing."""


@dataclass(frozen=True)
class TextExtractionResult:
    """Normalized text extracted from an evidence file."""

    text: str
    engine: str


class TextExtractor:
    """Extract textual content from various file formats using specialized libraries."""

    def __init__(self) -> None:
        self._pdfminer_available = None
        self._docx_available = None
        self._pptx_available = None
        self._xlsx_available = None
        self._beautifulsoup_available = None
        self._msg_available = None

    def extract(self, stream: BinaryIO, *, source_name: str) -> TextExtractionResult:
        """Return extracted text from the provided binary stream."""

        suffix = Path(source_name).suffix.lower() or ".bin"
        with self._materialize_stream(stream, suffix=suffix) as temp_path:
            engine = None
            text = ""

            # Route to appropriate extractor based on file type
            if suffix == ".pdf":
                text, engine = self._extract_pdf(temp_path)
            elif suffix in {".docx"}:
                text, engine = self._extract_docx(temp_path)
            elif suffix in {".doc"}:
                text, engine = self._extract_doc(temp_path)
            elif suffix in {".pptx"}:
                text, engine = self._extract_pptx(temp_path)
            elif suffix in {".xlsx", ".xls"}:
                text, engine = self._extract_xlsx(temp_path)
            elif suffix in {".xml", ".html", ".htm"}:
                text, engine = self._extract_xml_html(temp_path)
            elif suffix in {".eml"}:
                text, engine = self._extract_eml(temp_path)
            elif suffix in {".msg"}:
                text, engine = self._extract_msg(temp_path)
            elif suffix in {".json"}:
                text, engine = self._extract_json(temp_path)
            elif suffix in {".csv", ".tsv"}:
                text, engine = self._extract_csv(temp_path, suffix)
            elif suffix in {".txt", ".md", ".log", ".vcf", ".ics", ".rtf"}:
                text, engine = self._extract_plain_text(temp_path)
            else:
                # Try plain text as fallback
                text, engine = self._extract_plain_text(temp_path)

            return TextExtractionResult(text=text or "", engine=engine or "unknown")

    @contextmanager
    def _materialize_stream(self, stream: BinaryIO, *, suffix: str) -> Iterator[Path]:
        with NamedTemporaryFile(delete=False, suffix=suffix) as handle:
            temp_path = Path(handle.name)
            if hasattr(stream, "seek"):
                try:
                    stream.seek(0)
                except (OSError, AttributeError):
                    logger.debug("Stream for %s is not seekable; proceeding from current position", suffix)
            shutil.copyfileobj(stream, handle)

        try:
            yield temp_path
        finally:
            try:
                temp_path.unlink()
            except FileNotFoundError:  # pragma: no cover
                logger.debug("Temporary text file %s already removed", temp_path)

    def _extract_pdf(self, path: Path) -> tuple[str, str]:
        """Extract text from PDF files."""
        text = self._try_pdfminer(path)
        return (text, "pdfminer") if text else ("", "pdfminer")

    def _extract_docx(self, path: Path) -> tuple[str, str]:
        """Extract text from DOCX files."""
        try:
            import docx  # type: ignore[import]
        except ImportError:
            self._docx_available = False
            logger.debug("python-docx não disponível para %s", path.name)
            return ("", "docx-unavailable")
        
        try:
            doc = docx.Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return ("\n".join(paragraphs), "docx")
        except Exception as exc:
            logger.warning("docx falhou para %s: %s", path.name, exc)
            return ("", "docx-error")

    def _extract_doc(self, path: Path) -> tuple[str, str]:
        """Extract text from DOC files (legacy Word format)."""
        try:
            import docx2txt  # type: ignore[import]
        except ImportError:
            logger.debug("docx2txt não disponível para %s", path.name)
            return ("", "doc-unavailable")
        
        try:
            text = docx2txt.process(str(path))
            return (text or "", "docx2txt")
        except Exception as exc:
            logger.warning("docx2txt falhou para %s: %s", path.name, exc)
            return ("", "docx2txt-error")

    def _extract_pptx(self, path: Path) -> tuple[str, str]:
        """Extract text from PPTX files."""
        try:
            from pptx import Presentation  # type: ignore[import]
        except ImportError:
            self._pptx_available = False
            logger.debug("python-pptx não disponível para %s", path.name)
            return ("", "pptx-unavailable")
        
        try:
            prs = Presentation(str(path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text.strip())
            return ("\n".join(texts), "pptx")
        except Exception as exc:
            logger.warning("pptx falhou para %s: %s", path.name, exc)
            return ("", "pptx-error")

    def _extract_xlsx(self, path: Path) -> tuple[str, str]:
        """Extract text from XLSX/XLS files."""
        suffix = path.suffix.lower()
        try:
            if suffix == ".xlsx":
                import openpyxl  # type: ignore[import]
                wb = openpyxl.load_workbook(str(path), data_only=True)
                texts = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " ".join(str(cell) if cell is not None else "" for cell in row)
                        if row_text.strip():
                            texts.append(row_text)
                return ("\n".join(texts), "openpyxl")
            else:  # .xls
                import xlrd  # type: ignore[import]
                wb = xlrd.open_workbook(str(path))
                texts = []
                for sheet in wb.sheets():
                    for row_idx in range(sheet.nrows):
                        row = sheet.row_values(row_idx)
                        row_text = " ".join(str(cell) if cell else "" for cell in row)
                        if row_text.strip():
                            texts.append(row_text)
                return ("\n".join(texts), "xlrd")
        except ImportError:
            logger.debug("openpyxl/xlrd não disponível para %s", path.name)
            return ("", "xlsx-unavailable")
        except Exception as exc:
            logger.warning("xlsx falhou para %s: %s", path.name, exc)
            return ("", "xlsx-error")

    def _extract_xml_html(self, path: Path) -> tuple[str, str]:
        """Extract text from XML/HTML files."""
        try:
            from bs4 import BeautifulSoup  # type: ignore[import]
        except ImportError:
            self._beautifulsoup_available = False
            # Fallback to plain text extraction
            return self._extract_plain_text(path)
        
        try:
            with path.open("rb") as f:
                soup = BeautifulSoup(f, "html.parser" if path.suffix.lower() in {".html", ".htm"} else "xml")
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text(separator="\n", strip=True)
            return (text, "beautifulsoup")
        except Exception as exc:
            logger.warning("beautifulsoup falhou para %s: %s", path.name, exc)
            # Fallback to plain text
            return self._extract_plain_text(path)

    def _extract_eml(self, path: Path) -> tuple[str, str]:
        """Extract text from EML email files."""
        try:
            with path.open("rb") as f:
                msg = email.message_from_bytes(f.read())
            
            texts = []
            # Subject
            if msg["Subject"]:
                texts.append(f"Subject: {msg['Subject']}")
            
            # Body
            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    if content_type == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            texts.append(payload.decode("utf-8", errors="ignore"))
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    texts.append(payload.decode("utf-8", errors="ignore"))
            
            return ("\n".join(texts), "email")
        except Exception as exc:
            logger.warning("eml falhou para %s: %s", path.name, exc)
            return ("", "email-error")

    def _extract_msg(self, path: Path) -> tuple[str, str]:
        """Extract text from MSG email files."""
        try:
            import extract_msg  # type: ignore[import]
        except ImportError:
            self._msg_available = False
            logger.debug("extract-msg não disponível para %s", path.name)
            return ("", "msg-unavailable")
        
        try:
            msg = extract_msg.Message(str(path))
            texts = []
            if msg.subject:
                texts.append(f"Subject: {msg.subject}")
            if msg.body:
                texts.append(msg.body)
            return ("\n".join(texts), "extract-msg")
        except Exception as exc:
            logger.warning("extract-msg falhou para %s: %s", path.name, exc)
            return ("", "extract-msg-error")

    def _extract_json(self, path: Path) -> tuple[str, str]:
        """Extract text from JSON files."""
        try:
            with path.open("r", encoding="utf-8", errors="ignore") as f:
                data = json.load(f)
            # Convert JSON to readable text format
            text = json.dumps(data, indent=2, ensure_ascii=False)
            return (text, "json")
        except Exception as exc:
            logger.warning("json falhou para %s: %s", path.name, exc)
            return ("", "json-error")

    def _extract_csv(self, path: Path, suffix: str) -> tuple[str, str]:
        """Extract text from CSV/TSV files."""
        try:
            delimiter = "\t" if suffix == ".tsv" else ","
            with path.open("r", encoding="utf-8", errors="ignore") as f:
                reader = csv.reader(f, delimiter=delimiter)
                rows = []
                for row in reader:
                    rows.append(" | ".join(row))
            return ("\n".join(rows), "csv")
        except Exception as exc:
            logger.warning("csv falhou para %s: %s", path.name, exc)
            return ("", "csv-error")

    def _extract_plain_text(self, path: Path) -> tuple[str, str]:
        """Extract text from plain text files."""
        try:
            # Try UTF-8 first
            with path.open("r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            return (text, "plain-text")
        except Exception:
            # Fallback to latin-1
            try:
                with path.open("r", encoding="latin-1", errors="ignore") as f:
                    text = f.read()
                return (text, "plain-text-latin1")
            except Exception as exc:
                logger.warning("plain-text falhou para %s: %s", path.name, exc)
                return ("", "plain-text-error")

    def _try_pdfminer(self, path: Path) -> str:
        """Extract text from PDF using pdfminer."""
        if self._pdfminer_available is False:
            return ""

        try:
            from pdfminer.high_level import extract_text  # type: ignore[import]
        except ImportError:
            self._pdfminer_available = False
            logger.debug("pdfminer.six não disponível; retornando texto vazio para %s", path.name)
            return ""

        try:
            text = extract_text(str(path))
        except Exception as exc:  # pragma: no cover - defensivo
            logger.warning("pdfminer falhou para %s: %s", path, exc)
            self._pdfminer_available = False
            return ""

        self._pdfminer_available = True
        return text or ""
